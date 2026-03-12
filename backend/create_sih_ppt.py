from pptx import Presentation
from pptx.util import Inches, Pt

# Create a new presentation
prs = Presentation()

# Define slide layout: 0=Title Slide, 1=Title and Content
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

# 1. Title Slide
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "AI-Driven Crop Disease Prediction and Management System"
slide.placeholders[1].text = "Team Name: __________\nTeam Members: __________"

# 2. Problem Statement
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Problem Statement"
text_box = slide.placeholders[1]
text_box.text = ("Problem Statement Code & Title: AI-Driven Crop Disease Prediction and Management System\n\n"
                 "Background:\nCrop diseases can devastate yields, leading to significant financial losses for farmers. "
                 "Early detection and timely intervention are crucial for effective management.\n\n"
                 "Expected Solution:\nA mobile and web-based application that utilizes machine learning algorithms "
                 "to identify crop diseases and suggest preventive measures and treatments based on real-time data.")

# 3. Theme
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Theme"
slide.placeholders[1].text = "Agriculture / Smart Farming"

# 4. Abstract / Idea Summary
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Abstract / Idea Summary"
slide.placeholders[1].text = ("We propose an AI-driven system that analyzes crop images and environmental data "
                              "to predict potential disease outbreaks. The system provides actionable insights "
                              "and treatment recommendations to mitigate risks. Farmers can access this via "
                              "mobile and web apps, ensuring timely interventions and reducing financial losses.")

# 5. Innovation / Novelty
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Innovation / Novelty"
slide.placeholders[1].text = ("Unlike existing solutions that rely solely on manual observation or delayed "
                              "reporting, our system integrates real-time image analysis with environmental data "
                              "and provides immediate actionable advice. The model continuously learns and improves "
                              "from new data, making it highly adaptive and accurate.")

# 6. Proposed Solution
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Proposed Solution"
slide.placeholders[1].text = ("Technical Approach:\n- Machine Learning / Deep Learning for image-based disease detection\n"
                              "- Mobile and Web application interfaces\n"
                              "- Cloud backend for real-time data processing\n\n"
                              "Workflow:\n1. Farmer uploads crop image or sensor data\n2. System predicts disease\n"
                              "3. Provides treatment and preventive advice\n4. Stores data for continuous learning")

# 7. Usefulness / Relevance
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Usefulness / Relevance"
slide.placeholders[1].text = ("This system benefits farmers by enabling early detection of crop diseases, "
                              "reducing financial losses and improving crop yield. It is also useful for agricultural "
                              "extension officers to monitor and advise on disease management across regions.")

# 8. Feasibility & Scalability
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Feasibility & Scalability"
slide.placeholders[1].text = ("The solution is practical using existing smartphones, cameras, and cloud infrastructure. "
                              "It can be scaled to multiple crops and regions by retraining the AI model with new data.")

# 9. Expected Outcomes / Deliverables
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Expected Outcomes / Deliverables"
slide.placeholders[1].text = ("- AI-powered crop disease detection model\n"
                              "- Mobile application for farmers\n"
                              "- Web dashboard for monitoring and analytics\n"
                              "- Real-time treatment recommendations")

# 10. Potential Impact
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Potential Impact"
slide.placeholders[1].text = ("- Increased crop yield and reduced losses\n"
                              "- Economic benefits for farmers\n"
                              "- Environmental benefits through targeted treatment\n"
                              "- Scalable solution for large agricultural regions")

# 11. Tools / Technologies Required
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Tools / Technologies Required"
slide.placeholders[1].text = ("- Python, TensorFlow / Keras, OpenCV\n"
                              "- Flask / Django for web backend\n"
                              "- React Native / Flutter for mobile app\n"
                              "- Cloud platform: AWS / Google Cloud / Azure\n"
                              "- Database: MySQL / MongoDB")

# Save the presentation
prs.save("SIH_Crop_Disease_Abstract.pptx")
print("✅ PPT generated successfully: SIH_Crop_Disease_Abstract.pptx")
